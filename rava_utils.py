import itertools
import logging
import os
import shutil

import numpy as np
import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt, Inches
from re import sub

def find_raw_lyrics(project_path:str) -> list[str]:
    # Find the songs in tex folder.
    raw_folder: str = os.path.join(project_path, "lyrics", "00_raw")
    files: list[str] = [
        os.path.join(raw_folder, x)
        for x in os.listdir(raw_folder)
        if (x.endswith(".tex") or x.endswith('.txt'))
    ]

    return files

def remove_consecutive_blanks(l:list[str]) -> list[str]:
    #Remove consecutive '' elements in a list.
    res:list[str] = []

    for key, group in itertools.groupby(l, lambda x : x=='\n'):
        # check key=='' or not
        # if key!='' extend all values to res
        if not key:
            res.extend(list(group))
        # if key=='' append one '' to res
        else:
            res.append('\n')

    return res

def preprocess_tex(path_in: str, path_out: str, name: str) -> bool:
    """
    Preprocess the tex file.
    returns True when the file was valid
    returns False when the file was invalid
    
    lines are written to path_out, but if the file is invalid, the output file will not be touched
    """
    # Read the tex file.
    with open(path_in, "r") as f:
        tex_lines: np.ndarray = np.array(f.readlines())

    # Take everything between the \begin{obeylines} and \end{obeylines}.
    begin_obeylines: np.ndarray = np.argwhere(
            [
                r"\begin{obeylines}" in x or
                r"\begin{center}" in x 
                for x in tex_lines
            ]
    ).flatten()
    end_obeylines = np.argwhere([
        r"\end{obeylines}" in x or 
        r"\end{center}" in x 
        for x in tex_lines
        ]
    ).flatten()

    # Assert that we have the same number of \begin{obeylines} and \end{obeylines}.
    assert len(begin_obeylines) == len(
        end_obeylines
    ), f"Number of \\begin{{obeylines}} and \\end{{obeylines}} do not match in {path_in}."
    if len(begin_obeylines) == 0:
        return False

    clean_lines:list[str] = []

    # Loop over the \begin{obeylines} and \end{obeylines} and take everything in between.
    for idx1, idx2 in zip(begin_obeylines, end_obeylines):
        lines:list[str] = list(tex_lines[idx1 + 1 : idx2])
        lines = [x.replace("\\n", "\n").replace("\\newline", "\n").replace("\\\\", "\n") for x in lines]
        
        #new \vspace filtering
        for i in range(len(lines)):
            # lines[i] = sub(r"\vspace\{[0-9](.?[0-9]?){2}\}", "", lines[i])
            # lines[i] = sub(r"\vspace{.*?}", "", lines[i])
            lines[i] = sub(r"(-\ )?\\[a-zA-Z]*\{[a-zA-Z0-9.,\ ]*\}(:\ )?", "", lines[i])
            if lines[i].count("{") != lines[i].count("}"):
                lines[i] = lines[i].replace("{","").replace("}","")

        lines = [
            x
            for x in lines
            if not (
                x.startswith("\\") or x == "" or x.startswith("%")
            )
        ]

        clean_lines += lines

    clean_lines = remove_consecutive_blanks(clean_lines)

    if len(clean_lines) == 0:
        logging.info(f"No lines found in {path_in}. Skipping.")
        return False

    # Write the preprocessed lines to the outpath.
    with open(path_out, "w") as f:
        f.writelines(clean_lines)

    return True

def load_config(path:str) -> tuple[RGBColor, RGBColor, str, bool, Pt]:

    with open(path) as f:
        config = yaml.load(f, Loader=yaml.FullLoader)

    background_color = RGBColor(*tuple(config["project"]["background_color"]))
    font_color = RGBColor(*tuple(config["project"]["font_color"]))
    font_name = str(config["project"]["font_name"])
    font_bold = bool(config["project"]["font_bold"])
    font_size = Pt(config["project"]["font_size"])

    return background_color, font_color, font_name, font_bold, font_size

class PPTXSong:
    def __init__(self, background_color: RGBColor, font_color: RGBColor, font_name: str, font_bold: bool, font_size: Pt):

        self.prs = Presentation()
        self.background_color = background_color
        self.font_color = font_color
        self.font_name = font_name
        self.font_bold = font_bold
        self.font_size = font_size

    def add_slide(self, text:str):
        """
        Adds a slide with input string. \n means new line, and <blank> is a blank slide.
        """

        title_slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        if title is None:
            logging.warning("No title placeholder found in the slide layout. This should not happen, check the template.")
            return

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.background_color

        # Make the box take up the entire thing.
        title.left = Inches(0)
        slide_width:Emu | None = self.prs.slide_width # type: ignore
        slide_height:Emu | None = self.prs.slide_height # type: ignore
        if slide_width is None:
            logging.warning("No slide width found in the presentation. This should not happen, check the template.")
            return
        if slide_height is None:
            logging.warning("No slide height found in the presentation. This should not happen, check the template.")
            return
        title.width = slide_width
        title.height = slide_height

        # Add the text and change the font.
        text_frame = title.text_frame
        p = text_frame.paragraphs[0]

        n_line = 1
        for s in text.split("\\n"):
            # Remove leading and trailing spaces.
            s = s.strip()

            if n_line > 1:
                p.add_line_break()

            run = p.add_run()
            run.text = s

            # Add the formatting
            font = run.font
            font.name = self.font_name
            font.bold = self.font_bold
            font.size = self.font_size
            font.color.rgb = self.font_color

            # Signal we're jumping on to the next line.
            n_line += 1

        return

    def save(self, outpath:str):
        """
        Save the presentation.
        """
        self.prs.save(outpath)

def create_pptx(list_of_lyrics:list[str], out_path:str):

    # Load the config file.
    
    background_color: RGBColor
    font_color: RGBColor
    font_name: str
    font_bold: bool
    font_size: Pt

    background_color, font_color, font_name, font_bold, font_size = load_config(
        os.path.join(".", "revue_template", "config.yaml")
    )
    pptx = PPTXSong(background_color, font_color, font_name, font_bold, font_size)
    
    for line in list_of_lyrics:
        pptx.add_slide(line)

    pptx.save(out_path)
    return


def pptx_to_png(pptxPath:str, outfolder:str):
    assert pptxPath.endswith(".pptx"), "pptxPath should be a .pptx file"

    # Get the name of the song from the .pptx file.
    name = os.path.basename(pptxPath).replace(".pptx", "")

    if os.path.isdir(outfolder):
        # If the subfolder exists, nuke it.
        shutil.rmtree(outfolder)
        os.mkdir(outfolder)
    else:
        os.mkdir(outfolder)

    # Convert the pptx to a pdf using soffice.
    os.system(f"soffice --headless --convert-to pdf:writer_pdf_Export '{pptxPath}' > /dev/null 2>&1")

    # Grab the name of the infile (but as pdf).
    infile_pdf = pptxPath.replace(".pptx", ".pdf")
    
    #Move the file to where the pptx is located.
    os.rename(f'{name}.pdf', f'{infile_pdf}')

    # Convert the pdf to .pngs.
    os.system(
        f"convert -density 300 {infile_pdf} {outfolder}/{name}_%03d.png > /dev/null 2>&1"
    )

    # Delete the pdf in the /pptx/ folder afterwards.
    os.remove(infile_pdf)

    return
