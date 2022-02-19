#Dependencies
import os
import numpy as np
import argparse
import shutil

#Dependencies for doing PowerPoint magic.
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

class PPTXSong:
    def __init__(self, background_color = RGBColor(0,0,0), font_color = RGBColor(255,255,255), font_name  = 'Roboto', font_bold  = True, font_size  = Pt(45)):
        
        if background_color == 'black':
            background_color = RGBColor(0,0,0)
        elif background_color == 'white':
            background_color = RGBColor(1,1,1)
        
        
        self.prs = Presentation()
        self.background_color = background_color
        self.font_color       = font_color
        self.font_name        = font_name
        self.font_bold        = font_bold
        self.font_size        = font_size
        
    def add_slide(self, string):
        '''
        Adds a slide with input string. \n means new line, and <blank> is a blank slide.
        '''
        
        title_slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0,0,0)
    
        #Make the box take up the entire thing.
        title.left = 0
        title.width = self.prs.slide_width
        title.height = self.prs.slide_height

        #Add the text and change the font.
        text_frame = title.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = string
        
        #Add the formatting
        font = run.font
        font.name = self.font_name
        font.bold = self.font_bold
        font.size = self.font_size
        font.color.rgb = self.font_color

        return
    
    def save(self, outpath):
        '''
        Save the presentation.
        '''
        self.prs.save(outpath)

def tex_to_pptx(infile, outfile):
    '''
    Convert a .tex file to a .pptx file. 
    '''
    song = PPTXSong()

    #Remove leading and trailing spaces and double spaces too.
    lyrics = np.array([x.strip().replace('  ',' ') for x in open(infile, 'r').read().splitlines()])

    #Make sure that we have the same number of \begin{obeylines} and \end{obeylines}.
    assert sum([r'\begin{obeylines}' == x for x in lyrics]) == sum([r'\end{obeylines}' == x for x in lyrics]), r'Different number of \begin{obeylines} and \end{obeylines}.'

    lines = []

    #In case we have multiple \begin{obeylines} ... \end{obeylines}, then we loop over them. This cuts out the comments in between.
    for idx1, idx2 in zip(np.argwhere(lyrics==r'\begin{obeylines}').flatten(), np.argwhere(lyrics==r'\end{obeylines}').flatten()):
        lines = lyrics[idx1+1:idx2]

        #\n in the document is loaded as \\n. 
        lines = [x.replace('\\n', '\n') for x in lines if not (x.startswith('\\') or x == '' or x.startswith('%'))]

        for line in lines:
            if line == '<blank>': 
                line = ''
            song.add_slide(line)

    song.save(outfile)
    return

def download_from_overleaf(outfolder):
    '''
    Download the .tex files from /Musik/ folder in the Overleaf project to the local /tex/ folder.
    '''
    
    url = input('Enter the Overleaf URL: ')
    
    #Fetch the ID from the url.
    ID = url.split('/')[-1]
    
    os.system(f'git clone https://git.overleaf.com/{ID}/')
    
    assert os.path.isdir(f'./{ID}/Musik/'), 'The /Musik/ folder should be found in the overleaf project.'
    for filename in os.listdir(f'./{ID}/Musik/'):
        outfile = f'./{outfolder}/{filename}'
        
        #Move from the downloaded folder to the .tex folder.
        os.rename(f'./{ID}/Musik/{filename}', outfile)
    
    #Remove the folder afterwards.
    shutil.rmtree(f'./{ID}/')
    print('')

    return ID

def main():

    #Set up the arguments and parse it.
    parser = argparse.ArgumentParser()
    parser.add_argument('--folder', type=str, help='Path to the subfolder inside RAVA.')

    args = parser.parse_args()

    assert args.folder is not None, '--folder argument not given.'
    
    #Set the folder structure.
    if not os.path.isdir(args.folder):
        response = input('The folder does not exist. Do you want me to set it up? [y/n] ').lower().strip()
        if response == 'y': 
            #Make head folder
            os.mkdir(args.folder)
            #Make the subfolders. This should be changed if AV some day decide to change the file structure.
            for folder in ['lyrics', 'other', 'qlab', 'video', 'sound', 'pptx', 'tex']:
                os.mkdir(os.path.join(args.folder, folder))
        else:
            return
        print('Folder made.\n')

        #Download the songs from Overleaf.
        download_from_overleaf(outfolder = os.path.join(args.folder, 'tex'))
        print('/tex/ folder updated.')

    #Find the songs in tex folder.
    tex  = [os.path.join(args.folder, 'tex', x) for x in os.listdir(os.path.join(args.folder, 'tex')) if x.endswith('.tex')]

    for infile in tex:
       #Infile is .tex, outfile is .pptx
       outfile = infile.replace('/tex/', '/pptx/').replace('.tex','.pptx')
       
       #Convert the .tex to .pptx
       tex_to_pptx(infile, outfile)
    
    print('/pptx/ folder updated.')

    return


if __name__ == '__main__':
    main()

